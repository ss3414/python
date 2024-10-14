# ****************************************************************分割线****************************************************************
# todo 抽取文本

from pptx import Presentation

presentation = Presentation("C:/Users/Administrator/Desktop/test.pptx")
for slide in presentation.slides:
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                print(run.text)
